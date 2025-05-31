# app.py - Main Flask application for Vercel

import os
from flask import Flask, render_template, request, redirect, url_for, send_file, jsonify
import pandas as pd
import requests
from openai import OpenAI
from dotenv import load_dotenv
import tempfile
import yfinance as yf
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import io
import base64
from datetime import datetime, timedelta
import json

# Try to import matplotlib only if available
try:
    import matplotlib.pyplot as plt
    import matplotlib
    matplotlib.use('Agg')  # Use non-interactive backend for server environments
    MATPLOTLIB_AVAILABLE = True
except ImportError:
    MATPLOTLIB_AVAILABLE = False

# 1. Load environment variables from .env
load_dotenv()
NEWSAPI_KEY = os.getenv("NEWSAPI_KEY")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")

# 2. Initialize Flask
app = Flask(__name__)

# 3. Initialize OpenAI client
client = OpenAI(api_key=OPENAI_API_KEY)

# 4. Ensure required directories exist
upload_dir = "/tmp/uploads"
os.makedirs(upload_dir, exist_ok=True)
os.makedirs("templates", exist_ok=True)

# Health check endpoint for Vercel
@app.route("/health")
def health_check():
    return {"status": "healthy", "message": "Stock Exchange Email Generator is running"}, 200

def find_column(df, possible_names):
    """
    Helper function to find a column from a list of possible names
    """
    for name in possible_names:
        if name in df.columns:
            return name
    return None


# -----------------------------
# ROUTE #1: HOME / UPLOAD PAGE
# -----------------------------
@app.route("/", methods=["GET", "POST"])
def upload_file():
    if request.method == "POST":
        try:
            # A. Check if file part is in request
            if "file" not in request.files:
                return "No file part in the request", 400

            file = request.files["file"]
            if file.filename == "":
                return "No file selected", 400

            # B. Save the file to a temp location
            filepath = os.path.join("/tmp/uploads", file.filename)
            os.makedirs("/tmp/uploads", exist_ok=True)
            file.save(filepath)

            # C. Read the file with pandas (support CSV & Excel) - Enhanced with robust parsing
            try:
                if filepath.lower().endswith(".csv"):
                    # Try multiple approaches for robust CSV reading
                    try:
                        # First attempt: Standard reading
                        df = pd.read_csv(filepath)
                    except Exception as e1:
                        try:
                            # Second attempt: Handle malformed CSV with different options
                            df = pd.read_csv(filepath, 
                                           encoding='utf-8',
                                           on_bad_lines='skip')
                        except Exception as e2:
                            try:
                                # Third attempt: More lenient parsing
                                df = pd.read_csv(filepath, 
                                               sep=',',
                                               quotechar='"',
                                               skipinitialspace=True,
                                               on_bad_lines='skip',
                                               encoding='utf-8')
                            except Exception as e3:
                                try:
                                    # Fourth attempt: Try different encoding
                                    df = pd.read_csv(filepath, 
                                                   on_bad_lines='skip',
                                                   encoding='latin-1')
                                except Exception as e4:
                                    try:
                                        # Fifth attempt: Very lenient parsing for older pandas
                                        df = pd.read_csv(filepath, 
                                                       encoding='utf-8',
                                                       sep=None,
                                                       engine='python')
                                    except Exception as e5:
                                        return f"Error reading CSV file. Please check file format. Main error: {str(e1)[:150]}", 400
                else:
                    df = pd.read_excel(filepath)  # requires openpyxl installed
            except Exception as e:
                return f"Error reading file: {e}", 400

            # D. Intelligent column detection - find the main company identifier column
            company_column = None
            possible_company_names = ['Company', 'company', 'Company Name', 'company_name', 
                                    'CompanyName', 'Name', 'name', 'Symbol', 'Ticker', 
                                    'Company_Name', 'COMPANY', 'COMPANY_NAME', 'Issuer',
                                    'Company Legal Name', 'Legal Name']
            
            for col_name in possible_company_names:
                if col_name in df.columns:
                    company_column = col_name
                    break
            
            if company_column is None:
                # If no standard company column found, show available columns
                available_cols = ', '.join(df.columns.tolist())
                return f"Could not find a company name column. Available columns: {available_cols}. Please ensure your file has a column named 'Company', 'Name', or similar.", 400

            # E. Enhanced column standardization - map various column names to standard ones
            column_mappings = {
                # Company name (already handled above)
                'Company': company_column,
                
                # Exchange mappings
                'Exchange': find_column(df, ['Exchange', 'Stock_Exchange', 'Listed_Exchange', 'Market']),
                
                # Sector mappings
                'Sector': find_column(df, ['Sector', 'Primary Industry Sector', 'Industry_Sector', 'Business_Sector', 'Sector_Name']),
                
                # Industry mappings  
                'Industry': find_column(df, ['Industry', 'Primary Industry Group', 'Industry_Group', 'Sub_Industry', 'Industry_Name', 'Business_Industry']),
                
                # Country mappings
                'Country': find_column(df, ['Country', 'HQ Global Country/Territory', 'HQ_Country', 'Headquarters_Country', 'Location_Country']),
                
                # City mappings
                'City': find_column(df, ['City', 'HQ City', 'HQ_City', 'Headquarters_City', 'Location', 'Office_Location']),
                
                # Market cap mappings
                'Market_Cap_USD': find_column(df, ['Market_Cap_USD', 'Market Cap', 'Market_Cap', 'MarketCap', 'Market_Capitalization']),
                
                # Australian subsidiary mappings
                'Has_Australian_Subsidiary': find_column(df, ['Has_Australian_Subsidiary', 'Australian_Subsidiary', 'AUS_Subsidiary', 'Australia_Operations'])
            }
            
            # Apply column mappings (rename columns to standard names)
            rename_dict = {}
            for standard_name, source_column in column_mappings.items():
                if source_column and source_column in df.columns and source_column != standard_name:
                    rename_dict[source_column] = standard_name
            
            if rename_dict:
                df = df.rename(columns=rename_dict)
                print(f"Mapped columns: {rename_dict}")
                
                # Update company_column if it was renamed
                if company_column in rename_dict:
                    company_column = rename_dict[company_column]
            
            # F. Clean and validate the data
            # Remove rows where company name is missing
            df = df.dropna(subset=[company_column])
            df = df[df[company_column].astype(str).str.strip() != '']
            
            if len(df) == 0:
                return "No valid company data found in the file.", 400
            
            # Standardize the company column name to 'Company' for consistency
            if company_column != 'Company':
                df = df.rename(columns={company_column: 'Company'})
            
            # G. Show user what columns were detected
            detected_columns = df.columns.tolist()
            print(f"Successfully loaded {len(df)} companies with columns: {detected_columns}")
            
            # G. Store `df` in session‐like place with cleaned data
            df.to_csv("/tmp/uploads/last_upload.csv", index=False)

            # F. Redirect to the "select companies" page
            return redirect(url_for("select_companies"))
            
        except Exception as e:
            # Catch any unexpected errors and return a helpful message
            import traceback
            error_details = str(e)
            print(f"Upload error: {error_details}")
            print(f"Traceback: {traceback.format_exc()}")
            return f"An error occurred while processing your file: {error_details}. Please try again or check your file format.", 500

    # If GET, just render the upload form
    return render_template("upload.html")


# ----------------------------------
# ROUTE #2: SELECT COMPANIES TO PROCESS
# ----------------------------------
@app.route("/select", methods=["GET", "POST"])
def select_companies():
    # A. Load the stored DataFrame
    try:
        df = pd.read_csv("/tmp/uploads/last_upload.csv")
    except FileNotFoundError:
        return redirect(url_for("upload_file"))

    # B. Apply filters if provided - Enhanced to handle missing columns
    filtered_df = df.copy()
    
    # Helper function to safely get column values
    def safe_filter(dataframe, column_name, filter_value):
        if column_name in dataframe.columns and filter_value:
            return dataframe[dataframe[column_name] == filter_value]
        return dataframe
    
    def safe_filter_boolean(dataframe, column_name, filter_value):
        if column_name in dataframe.columns and filter_value:
            if filter_value == 'Yes':
                return dataframe[dataframe[column_name] == 'Yes']
            elif filter_value == 'No':
                return dataframe[dataframe[column_name] == 'No']
        return dataframe
    
    def safe_filter_market_cap(dataframe, filter_value):
        # Find market cap column
        market_cap_cols = ['Market_Cap_USD', 'Market_Cap', 'MarketCap', 'Market_Capitalization']
        market_cap_col = None
        for col in market_cap_cols:
            if col in dataframe.columns:
                market_cap_col = col
                break
        
        if not market_cap_col or not filter_value:
            return dataframe
        
        try:
            # Convert market cap to numeric values (handle various formats)
            def parse_market_cap(value):
                if pd.isna(value):
                    return 0
                value_str = str(value).replace(',', '').replace('$', '').strip()
                if 'B' in value_str.upper():
                    return float(value_str.upper().replace('B', '')) * 1000000000
                elif 'M' in value_str.upper():
                    return float(value_str.upper().replace('M', '')) * 1000000
                else:
                    try:
                        return float(value_str)
                    except:
                        return 0
            
            dataframe['_market_cap_numeric'] = dataframe[market_cap_col].apply(parse_market_cap)
            
            if filter_value == 'under_100m':
                return dataframe[dataframe['_market_cap_numeric'] < 100000000]
            elif filter_value == '100m_500m':
                return dataframe[(dataframe['_market_cap_numeric'] >= 100000000) & (dataframe['_market_cap_numeric'] < 500000000)]
            elif filter_value == '500m_1b':
                return dataframe[(dataframe['_market_cap_numeric'] >= 500000000) & (dataframe['_market_cap_numeric'] < 1000000000)]
            elif filter_value == '1b_5b':
                return dataframe[(dataframe['_market_cap_numeric'] >= 1000000000) & (dataframe['_market_cap_numeric'] < 5000000000)]
            elif filter_value == 'over_5b':
                return dataframe[dataframe['_market_cap_numeric'] >= 5000000000]
            
        except Exception as e:
            print(f"Market cap filtering error: {e}")
            return dataframe
        
        return dataframe
    
    def safe_filter_location(dataframe, filter_value):
        # Find location columns
        location_cols = ['City', 'Location', 'Headquarters', 'Australian_Office_Location', 'Office_Location']
        
        if not filter_value:
            return dataframe
        
        location_mask = pd.Series([False] * len(dataframe), index=dataframe.index)
        
        for col in location_cols:
            if col in dataframe.columns:
                # Case-insensitive partial matching
                mask = dataframe[col].astype(str).str.contains(filter_value, case=False, na=False)
                location_mask = location_mask | mask
        
        return dataframe[location_mask] if location_mask.any() else dataframe
    
    # Filter options with safe checking
    filter_exchange = request.args.get('exchange', '')
    filter_country = request.args.get('country', '')
    filter_australian_sub = request.args.get('australian_sub', '')
    filter_sector = request.args.get('sector', '')
    filter_industry = request.args.get('industry', '')  # New industry filter
    filter_market_cap = request.args.get('market_cap_range', '')  # New market cap filter
    filter_location = request.args.get('location', '')  # New location filter
    
    # Apply filters only if columns exist
    filtered_df = safe_filter(filtered_df, 'Exchange', filter_exchange)
    filtered_df = safe_filter(filtered_df, 'Country', filter_country)
    filtered_df = safe_filter_boolean(filtered_df, 'Has_Australian_Subsidiary', filter_australian_sub)
    filtered_df = safe_filter(filtered_df, 'Sector', filter_sector)
    filtered_df = safe_filter(filtered_df, 'Industry', filter_industry)  # Industry filter
    filtered_df = safe_filter_market_cap(filtered_df, filter_market_cap)  # Market cap filter
    filtered_df = safe_filter_location(filtered_df, filter_location)  # Location filter

    companies = filtered_df["Company"].tolist()
    companies_with_details = []
    
    # Get additional details for display - handle missing columns gracefully
    for idx, row in filtered_df.iterrows():
        company_info = {
            'original_index': df.index[df['Company'] == row['Company']].tolist()[0],
            'name': row['Company'],
            'exchange': row.get('Exchange', 'N/A'),
            'country': row.get('Country', 'N/A'),
            'sector': row.get('Sector', 'N/A'),
            'industry': row.get('Industry', 'N/A'),
            'city': row.get('City', row.get('Location', 'N/A')),
            'market_cap': row.get('Market_Cap_USD', row.get('Market_Cap', row.get('MarketCap', 'N/A'))),
            'has_aus_sub': row.get('Has_Australian_Subsidiary', row.get('Australian_Subsidiary', 'N/A'))
        }
        companies_with_details.append(company_info)

    # Get unique values for filter dropdowns - only for existing columns
    exchanges = sorted([x for x in df['Exchange'].unique() if pd.notna(x)]) if 'Exchange' in df.columns else []
    countries = sorted([x for x in df['Country'].unique() if pd.notna(x)]) if 'Country' in df.columns else []
    sectors = sorted([x for x in df['Sector'].unique() if pd.notna(x)]) if 'Sector' in df.columns else []
    industries = sorted([x for x in df['Industry'].unique() if pd.notna(x)]) if 'Industry' in df.columns else []
    
    # Generate dynamic location options from multiple columns
    locations = set()
    location_cols = ['City', 'Location', 'Headquarters', 'Australian_Office_Location', 'Office_Location']
    for col in location_cols:
        if col in df.columns:
            locations.update([x for x in df[col].dropna().unique() if str(x).strip() and str(x) != 'N/A'])
    locations = sorted(list(locations))
    
    # Generate market cap ranges based on actual data
    market_cap_ranges = [
        ('under_100m', 'Under $100M'),
        ('100m_500m', '$100M - $500M'), 
        ('500m_1b', '$500M - $1B'),
        ('1b_5b', '$1B - $5B'),
        ('over_5b', 'Over $5B')
    ]

    if request.method == "POST":
        # B. Get list of selected company indices
        selected = request.form.getlist("company")  # e.g. ["0", "3", "5"]
        if not selected:
            return "No companies selected", 400

        # Save the indices into a temp file for step 3
        with open("/tmp/uploads/selected_companies.txt", "w") as f:
            f.write(",".join(selected))
        return redirect(url_for("generate_emails"))

    # If GET, render a page with checkboxes and filters
    return render_template("select.html", 
                         companies=companies_with_details,
                         exchanges=exchanges,
                         countries=countries,
                         sectors=sectors,
                         industries=industries,
                         locations=locations,
                         market_cap_ranges=market_cap_ranges,
                         current_filters={
                             'exchange': filter_exchange,
                             'country': filter_country,
                             'australian_sub': filter_australian_sub,
                             'sector': filter_sector,
                             'industry': filter_industry,
                             'market_cap_range': filter_market_cap,
                             'location': filter_location
                         })


# ----------------------------
# ROUTE #3: GENERATE EMAILS
# ----------------------------
@app.route("/generate", methods=["GET"])
def generate_emails():
    # A. Load DataFrame and selected indices
    df = pd.read_csv("/tmp/uploads/last_upload.csv")
    try:
        with open("/tmp/uploads/selected_companies.txt", "r") as f:
            selected_indices = [int(x) for x in f.read().split(",") if x.strip().isdigit()]
    except FileNotFoundError:
        return redirect(url_for("select_companies"))

    results = []  # list of dicts: { company, market_cap, headlines, email_draft, pptx_path, company_data }

    for idx in selected_indices:
        if idx < 0 or idx >= len(df):
            continue
        row = df.iloc[idx]
        company_name = str(row["Company"])
        
        # Extract comprehensive company information - handle missing columns gracefully
        market_cap = row.get("Market_Cap_USD", row.get("Market_Cap", row.get("MarketCap", "Unknown")))
        exchange = row.get("Exchange", row.get("Stock_Exchange", "Unknown"))
        sector = row.get("Sector", row.get("Industry", row.get("Business_Sector", "Unknown")))
        country = row.get("Country", row.get("Headquarters", row.get("Location", "Unknown")))
        has_aus_sub = row.get("Has_Australian_Subsidiary", row.get("Australian_Subsidiary", row.get("AUS_Subsidiary", "Unknown")))
        ceo_name = row.get("CEO_Name", row.get("CEO", row.get("Chief_Executive", "Unknown")))
        ceo_email = row.get("CEO_Email", row.get("CEO_Contact", "Unknown"))
        ir_contact = row.get("IR_Contact_Email", row.get("IR_Email", row.get("Investor_Relations", "Unknown")))

        # Try to extract ticker symbol for financial data lookup
        ticker_symbol = extract_ticker_symbol(company_name, exchange)
        print(f"Processing {company_name}, ticker: {ticker_symbol}")

        # Fetch comprehensive company data for PowerPoint
        comprehensive_data = fetch_comprehensive_company_data(company_name, ticker_symbol)

        # 1. Fetch recent news from NewsAPI
        news_url = (
            f"https://newsapi.org/v2/everything?"
            f"q={requests.utils.quote(company_name)}&"
            f"sortBy=publishedAt&"
            f"pageSize=5&"              # grab top 5 articles for more content
            f"apiKey={NEWSAPI_KEY}"
        )
        news_resp = requests.get(news_url)
        news_data = news_resp.json()
        headlines = []
        if news_data.get("status") == "ok":
            # Extract up to 5 headlines with descriptions
            for article in news_data.get("articles", [])[:5]:
                title = article.get("title", "")
                description = article.get("description", "")
                url = article.get("url", "")
                published_at = article.get("publishedAt", "")
                
                if title:
                    # Create rich content for better email personalization
                    if description and len(description) > 20:
                        headlines.append(f"• {title}\n  Summary: {description[:150]}{'...' if len(description) > 150 else ''}")
                    else:
                        headlines.append(f"• {title}")
        
        # Add fallback content if no news found
        if not headlines:
            headlines.append(f"Recent market activities and business developments for {company_name} indicate continued growth potential in their {sector.lower()} sector.")
        
        # Limit to top 3 most relevant headlines for the prompt
        headlines = headlines[:3]

        # 2. Call OpenAI to generate a personalized email draft
        # Build a more comprehensive prompt
        company_data = analyze_company_data(row, df.columns)
        prompt = create_intelligent_prompt(company_name, company_data, headlines)
        try:
            completion = client.chat.completions.create(
                model="gpt-3.5-turbo",  # free-trial tier model
                messages=[
                    {"role": "system", "content": "You are an email generator specializing in stock exchange listing outreach."},
                    {"role": "user", "content": prompt}
                ],
                max_tokens=250,
                temperature=0.7,
            )
            email_text = completion.choices[0].message.content.strip()
        except Exception as e:
            email_text = f"Error generating email: {e}"

        # 3. Generate PowerPoint presentation
        pptx_path = None
        try:
            pptx_path = create_company_powerpoint(company_name, comprehensive_data, row)
            print(f"PowerPoint created for {company_name}: {pptx_path}")
        except Exception as e:
            print(f"Error creating PowerPoint for {company_name}: {e}")

        results.append({
            "company": company_name,
            "exchange": exchange,
            "sector": sector,
            "country": country,
            "market_cap": market_cap,
            "has_aus_sub": has_aus_sub,
            "ceo_name": ceo_name,
            "ceo_email": ceo_email,
            "ir_contact": ir_contact,
            "headlines": headlines,
            "email_draft": email_text,
            "analyzed_data": company_data,  # Include the analyzed data for transparency
            "pptx_path": pptx_path,  # PowerPoint file path
            "comprehensive_data": comprehensive_data,  # Full financial data
            "ticker_symbol": ticker_symbol  # Ticker symbol used
        })

    # 3. Render results in an HTML page
    return render_template("results.html", results=results)

# Route to download PowerPoint files
@app.route("/download_pptx/<filename>")
def download_pptx(filename):
    """
    Download PowerPoint presentations
    """
    try:
        file_path = f"/tmp/uploads/{filename}"
        if os.path.exists(file_path):
            return send_file(file_path, as_attachment=True, download_name=filename)
        else:
            return "File not found", 404
    except Exception as e:
        return f"Error downloading file: {e}", 500

def analyze_company_data(row, df_columns):
    """
    Intelligently analyze all available company data and categorize by relevance to stock exchange listing
    """
    company_data = {
        'financial_metrics': {},
        'business_info': {},
        'geographic_info': {},
        'leadership_info': {},
        'recent_activities': {},
        'other_relevant': {}
    }
    
    # Define column categories for stock exchange listing relevance
    financial_keywords = ['market_cap', 'revenue', 'profit', 'ebitda', 'assets', 'debt', 'cash', 'capital', 'valuation', 'shares', 'price', 'dividend']
    business_keywords = ['sector', 'industry', 'business', 'description', 'products', 'services', 'operations', 'subsidiaries']
    geographic_keywords = ['country', 'headquarters', 'location', 'region', 'australian', 'asia', 'pacific', 'office']
    leadership_keywords = ['ceo', 'cfo', 'chairman', 'director', 'executive', 'management', 'contact', 'email', 'ir', 'investor']
    activity_keywords = ['ipo', 'listing', 'raising', 'round', 'investment', 'acquisition', 'merger', 'expansion', 'growth']
    
    for column in df_columns:
        if column == 'Company':
            continue
            
        value = row.get(column, '')
        if pd.isna(value) or str(value).strip() == '' or str(value).lower() in ['unknown', 'n/a', 'na', 'none']:
            continue
            
        column_lower = column.lower()
        value_str = str(value)
        
        # Categorize by column name and content
        if any(keyword in column_lower for keyword in financial_keywords):
            company_data['financial_metrics'][column] = value_str
        elif any(keyword in column_lower for keyword in business_keywords):
            company_data['business_info'][column] = value_str
        elif any(keyword in column_lower for keyword in geographic_keywords):
            company_data['geographic_info'][column] = value_str
        elif any(keyword in column_lower for keyword in leadership_keywords):
            company_data['leadership_info'][column] = value_str
        elif any(keyword in column_lower for keyword in activity_keywords):
            company_data['recent_activities'][column] = value_str
        else:
            # Check if the content might be relevant
            if len(value_str) > 10 and any(word in value_str.lower() for word in ['million', 'billion', 'usd', 'aud', 'growth', 'market', 'international']):
                company_data['other_relevant'][column] = value_str
    
    return company_data

def create_intelligent_prompt(company_name, company_data, headlines):
    """
    Create a highly personalized prompt based on all available company data
    """
    
    # Extract key information
    financial_info = company_data['financial_metrics']
    business_info = company_data['business_info']
    geographic_info = company_data['geographic_info']
    leadership_info = company_data['leadership_info']
    activities_info = company_data['recent_activities']
    other_info = company_data['other_relevant']
    
    # Build dynamic sections based on available data
    prompt_sections = []
    
    prompt_sections.append(f"You are writing a highly personalized stock exchange listing proposal email to {company_name}.")
    
    # Company Profile Section
    profile_details = []
    if business_info:
        for key, value in business_info.items():
            if len(value) < 100:  # Avoid overly long descriptions
                profile_details.append(f"{key}: {value}")
    
    if financial_info:
        for key, value in financial_info.items():
            profile_details.append(f"{key}: {value}")
    
    if geographic_info:
        for key, value in geographic_info.items():
            profile_details.append(f"{key}: {value}")
    
    if profile_details:
        prompt_sections.append(f"COMPANY PROFILE:\n" + "\n".join(f"- {detail}" for detail in profile_details))
    
    # Leadership Information
    if leadership_info:
        leadership_details = []
        for key, value in leadership_info.items():
            leadership_details.append(f"{key}: {value}")
        prompt_sections.append(f"LEADERSHIP:\n" + "\n".join(f"- {detail}" for detail in leadership_details))
    
    # Recent Activities
    if activities_info:
        activity_details = []
        for key, value in activities_info.items():
            activity_details.append(f"{key}: {value}")
        prompt_sections.append(f"RECENT CORPORATE ACTIVITIES:\n" + "\n".join(f"- {detail}" for detail in activity_details))
    
    # News Analysis
    if headlines:
        prompt_sections.append(f"RECENT NEWS ANALYSIS:\n" + "\n".join(headlines))
    
    # Strategic recommendations based on available data
    strategic_points = []
    
    # Australian connection
    aus_connection = any('australian' in str(v).lower() for v in geographic_info.values())
    if aus_connection:
        strategic_points.append("- Leverage existing Australian presence for natural market progression")
    
    # Financial scale
    has_large_market_cap = any('billion' in str(v).lower() or 'million' in str(v).lower() for v in financial_info.values())
    if has_large_market_cap:
        strategic_points.append("- Scale and financial profile suitable for institutional investors")
    
    # Geographic expansion
    asia_pacific_presence = any(region in str(geographic_info).lower() for region in ['asia', 'pacific', 'singapore', 'hong kong', 'japan'])
    if asia_pacific_presence:
        strategic_points.append("- Asia-Pacific presence aligns with regional investor base")
    
    # Add other relevant information
    if other_info:
        other_details = []
        for key, value in other_info.items():
            if len(value) < 150:  # Keep it concise
                other_details.append(f"{key}: {value}")
        if other_details:
            prompt_sections.append(f"ADDITIONAL RELEVANT INFORMATION:\n" + "\n".join(f"- {detail}" for detail in other_details))
    
    # Final instructions
    prompt_sections.append("""
TASK: Write a compelling, research-driven email that:
1. Opens with specific reference to their recent developments or business profile
2. Demonstrates deep understanding of their business and current situation
3. Connects their specific circumstances to stock exchange listing advantages
4. Creates urgency based on their growth trajectory and market position
5. Includes a clear call-to-action

STRATEGIC FOCUS AREAS:""")
    
    if strategic_points:
        prompt_sections.append("\n".join(strategic_points))
    else:
        prompt_sections.append("- Access to international capital markets\n- Regulatory advantages and investor familiarity\n- Currency diversification benefits")
    
    prompt_sections.append("""
TONE: Professional, well-researched, and consultative (demonstrate you've done homework)
LENGTH: 175-225 words
FORMAT: Professional business email with clear next steps
REQUIREMENT: Make it clear this email is based on thorough research of their specific situation.""")
    
    return "\n\n".join(prompt_sections)

# PowerPoint Generation Functions

def fetch_comprehensive_company_data(company_name, ticker_symbol=None):
    """
    Fetch comprehensive company data from multiple sources including financial metrics,
    stock price data, and company information
    """
    company_data = {
        'basic_info': {},
        'financial_metrics': {},
        'stock_data': {},
        'price_chart_data': None,
        'key_stats': {},
        'major_shareholders': [],
        'recent_news': [],
        'capital_raising': [],
        'error': None
    }
    
    try:
        # Try to fetch data using yfinance if we have a ticker
        if ticker_symbol:
            try:
                stock = yf.Ticker(ticker_symbol)
                
                # Get basic company info
                info = stock.info
                company_data['basic_info'] = {
                    'name': info.get('longName', company_name),
                    'sector': info.get('sector', 'N/A'),
                    'industry': info.get('industry', 'N/A'),
                    'country': info.get('country', 'N/A'),
                    'website': info.get('website', 'N/A'),
                    'description': info.get('longBusinessSummary', 'N/A')[:500] + '...' if info.get('longBusinessSummary') else 'N/A'
                }
                
                # Get financial metrics
                company_data['financial_metrics'] = {
                    'market_cap': info.get('marketCap'),
                    'revenue': info.get('totalRevenue'),
                    'profit_margin': info.get('profitMargins'),
                    'pe_ratio': info.get('trailingPE'),
                    'price_to_book': info.get('priceToBook'),
                    'debt_to_equity': info.get('debtToEquity'),
                    'return_on_equity': info.get('returnOnEquity'),
                    'revenue_growth': info.get('revenueGrowth')
                }
                
                # Get current stock data
                company_data['stock_data'] = {
                    'current_price': info.get('currentPrice'),
                    'previous_close': info.get('previousClose'),
                    'day_change': info.get('currentPrice', 0) - info.get('previousClose', 0) if info.get('currentPrice') and info.get('previousClose') else 0,
                    'day_change_percent': ((info.get('currentPrice', 0) - info.get('previousClose', 0)) / info.get('previousClose', 1)) * 100 if info.get('currentPrice') and info.get('previousClose') else 0,
                    '52_week_high': info.get('fiftyTwoWeekHigh'),
                    '52_week_low': info.get('fiftyTwoWeekLow'),
                    'volume': info.get('volume'),
                    'avg_volume': info.get('averageVolume')
                }
                
                # Get key statistics
                company_data['key_stats'] = {
                    'employees': info.get('fullTimeEmployees'),
                    'dividend_yield': info.get('dividendYield'),
                    'payout_ratio': info.get('payoutRatio'),
                    'beta': info.get('beta'),
                    'shares_outstanding': info.get('sharesOutstanding'),
                    'float_shares': info.get('floatShares')
                }
                
                # Get major shareholders (top 10)
                try:
                    major_holders = stock.major_holders
                    if major_holders is not None and not major_holders.empty:
                        company_data['major_shareholders'] = major_holders.to_dict('records')[:5]  # Top 5
                except:
                    pass
                
                # Get 5-year price data for chart
                try:
                    end_date = datetime.now()
                    start_date = end_date - timedelta(days=5*365)  # 5 years
                    hist_data = stock.history(start=start_date, end=end_date, interval='1mo')  # Monthly data
                    
                    if not hist_data.empty:
                        company_data['price_chart_data'] = {
                            'dates': [date.strftime('%Y-%m') for date in hist_data.index],
                            'prices': hist_data['Close'].tolist(),
                            'volumes': hist_data['Volume'].tolist()
                        }
                except Exception as e:
                    print(f"Error fetching price data: {e}")
                
                # Get recent news
                try:
                    news = stock.news[:5]  # Get top 5 news items
                    company_data['recent_news'] = [
                        {
                            'title': item.get('title', ''),
                            'publisher': item.get('publisher', ''),
                            'link': item.get('link', ''),
                            'published': datetime.fromtimestamp(item.get('providerPublishTime', 0)).strftime('%Y-%m-%d') if item.get('providerPublishTime') else ''
                        }
                        for item in news if item.get('title')
                    ]
                except:
                    pass
                    
            except Exception as e:
                company_data['error'] = f"Error fetching data for {ticker_symbol}: {str(e)}"
                print(f"Error fetching yfinance data: {e}")
        
        # If no ticker or yfinance failed, try to search for more info
        if not company_data['basic_info'] or company_data['error']:
            company_data['basic_info'] = {
                'name': company_name,
                'sector': 'N/A',
                'industry': 'N/A', 
                'country': 'N/A',
                'website': 'N/A',
                'description': f'No detailed information available for {company_name}. This company may be privately held or listed on a regional exchange.'
            }
    
    except Exception as e:
        company_data['error'] = f"General error: {str(e)}"
        print(f"General error in fetch_comprehensive_company_data: {e}")
    
    return company_data

def create_stock_price_chart(company_data, company_name):
    """
    Create a 5-year stock price chart using matplotlib if available,
    otherwise create a simple text representation
    """
    try:
        if not MATPLOTLIB_AVAILABLE:
            # Create a simple text-based chart representation
            if company_data.get('price_chart_data'):
                price_data = company_data['price_chart_data']
                prices = price_data['prices']
                if prices:
                    min_price = min(prices)
                    max_price = max(prices)
                    latest_price = prices[-1]
                    change_5y = ((latest_price - prices[0]) / prices[0]) * 100 if prices[0] > 0 else 0
                    
                    chart_text = f"""
5-Year Stock Performance Summary for {company_name}:
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━

Starting Price: ${prices[0]:.2f}
Current Price:  ${latest_price:.2f}
5-Year Change:  {change_5y:+.1f}%

Price Range:
Highest: ${max_price:.2f}
Lowest:  ${min_price:.2f}

Note: Detailed chart available in downloadable PowerPoint presentation
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
                    """
                    return chart_text
            return f"Stock price data not available for {company_name}"
        
        if not company_data.get('price_chart_data'):
            # Create a placeholder chart
            fig, ax = plt.subplots(figsize=(10, 6))
            ax.text(0.5, 0.5, f'Stock price data not available for {company_name}', 
                   horizontalalignment='center', verticalalignment='center', 
                   transform=ax.transAxes, fontsize=14)
            ax.set_title(f'{company_name} - 5 Year Stock Price Chart', fontsize=16, fontweight='bold')
            plt.tight_layout()
        else:
            # Create actual chart
            price_data = company_data['price_chart_data']
            dates = price_data['dates']
            prices = price_data['prices']
            
            fig, ax = plt.subplots(figsize=(12, 7))
            
            # Plot the price line
            ax.plot(dates, prices, linewidth=2.5, color='#2E86C1', alpha=0.8)
            ax.fill_between(dates, prices, alpha=0.3, color='#85C1E9')
            
            # Styling
            ax.set_title(f'{company_name} - 5 Year Stock Price Performance', 
                        fontsize=16, fontweight='bold', pad=20)
            ax.set_xlabel('Year', fontsize=12)
            ax.set_ylabel('Stock Price ($)', fontsize=12)
            ax.grid(True, alpha=0.3)
            
            # Format x-axis to show fewer labels
            ax.set_xticks(ax.get_xticks()[::6])  # Show every 6th tick
            
            # Add some styling
            ax.spines['top'].set_visible(False)
            ax.spines['right'].set_visible(False)
            
            plt.xticks(rotation=45)
            plt.tight_layout()
        
        # Save to bytes
        img_buffer = io.BytesIO()
        plt.savefig(img_buffer, format='png', dpi=300, bbox_inches='tight')
        img_buffer.seek(0)
        
        # Convert to base64
        img_base64 = base64.b64encode(img_buffer.getvalue()).decode()
        plt.close()
        
        return img_base64
        
    except Exception as e:
        print(f"Error creating chart: {e}")
        return None

def create_company_powerpoint(company_name, company_data, financial_data_from_csv=None):
    """
    Create a comprehensive PowerPoint presentation for the company
    """
    try:
        # Create presentation
        prs = Presentation()
        
        # Set slide dimensions (16:9 aspect ratio)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        
        # Slide 1: Company Overview
        slide1 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        
        # Title
        title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = f"{company_data['basic_info']['name']} - Investment Overview"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(44, 62, 80)
        title_para.alignment = PP_ALIGN.CENTER
        
        # Company description
        desc_box = slide1.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(3))
        desc_frame = desc_box.text_frame
        desc_frame.word_wrap = True
        
        # Basic company info
        basic_info = company_data['basic_info']
        description = f"""
COMPANY OVERVIEW:
{basic_info.get('description', 'N/A')}

SECTOR: {basic_info.get('sector', 'N/A')}
INDUSTRY: {basic_info.get('industry', 'N/A')}
COUNTRY: {basic_info.get('country', 'N/A')}
        """
        
        desc_frame.text = description.strip()
        for paragraph in desc_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = RGBColor(52, 73, 94)
        
        # Key metrics box
        metrics_box = slide1.shapes.add_textbox(Inches(1), Inches(5), Inches(11), Inches(2))
        metrics_frame = metrics_box.text_frame
        
        financial_metrics = company_data['financial_metrics']
        stock_data = company_data['stock_data']
        
        def format_number(value, is_currency=False, is_percentage=False):
            if value is None or value == 'N/A':
                return 'N/A'
            try:
                if is_percentage:
                    return f"{float(value)*100:.1f}%" if value < 1 else f"{float(value):.1f}%"
                elif is_currency:
                    if float(value) >= 1e9:
                        return f"${float(value)/1e9:.1f}B"
                    elif float(value) >= 1e6:
                        return f"${float(value)/1e6:.1f}M"
                    else:
                        return f"${float(value):,.0f}"
                else:
                    return f"{float(value):,.2f}"
            except:
                return str(value)
        
        metrics_text = f"""KEY FINANCIAL METRICS:
Market Cap: {format_number(financial_metrics.get('market_cap'), is_currency=True)}  |  Current Price: {format_number(stock_data.get('current_price'), is_currency=True)}  |  P/E Ratio: {format_number(financial_metrics.get('pe_ratio'))}
Revenue: {format_number(financial_metrics.get('revenue'), is_currency=True)}  |  Profit Margin: {format_number(financial_metrics.get('profit_margin'), is_percentage=True)}  |  Beta: {format_number(company_data['key_stats'].get('beta'))}"""
        
        metrics_frame.text = metrics_text
        for paragraph in metrics_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(39, 174, 96)
        
        # Slide 2: Financial Performance & Stock Chart
        slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        
        # Title
        title_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.8))
        title_frame2 = title_box2.text_frame
        title_frame2.text = f"{company_name} - Financial Performance & Market Data"
        title_para2 = title_frame2.paragraphs[0]
        title_para2.font.size = Pt(24)
        title_para2.font.bold = True
        title_para2.font.color.rgb = RGBColor(44, 62, 80)
        title_para2.alignment = PP_ALIGN.CENTER
        
        # Create and add stock chart
        chart_base64 = create_stock_price_chart(company_data, company_name)
        if chart_base64:
            if MATPLOTLIB_AVAILABLE and not isinstance(chart_base64, str):
                # Save chart to temp file (base64 image)
                chart_data = base64.b64decode(chart_base64)
                chart_path = f"/tmp/uploads/{company_name}_chart.png"
                with open(chart_path, 'wb') as f:
                    f.write(chart_data)
                
                # Add chart to slide
                slide2.shapes.add_picture(chart_path, Inches(0.5), Inches(1.2), width=Inches(7))
            else:
                # Add text-based chart (when matplotlib not available)
                chart_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(7), Inches(4))
                chart_frame = chart_box.text_frame
                chart_frame.word_wrap = True
                chart_frame.text = chart_base64 if isinstance(chart_base64, str) else f"Chart data not available for {company_name}"
                
                for paragraph in chart_frame.paragraphs:
                    paragraph.font.size = Pt(10)
                    paragraph.font.name = 'Courier New'  # Monospace font for better text chart display
                    paragraph.font.color.rgb = RGBColor(52, 73, 94)
        
        # Financial summary on the right
        summary_box = slide2.shapes.add_textbox(Inches(8), Inches(1.2), Inches(5), Inches(5.5))
        summary_frame = summary_box.text_frame
        summary_frame.word_wrap = True
        
        # Create financial summary
        summary_text = f"""INVESTMENT HIGHLIGHTS:

📈 STOCK PERFORMANCE:
• Current Price: {format_number(stock_data.get('current_price'), is_currency=True)}
• 52W High: {format_number(stock_data.get('52_week_high'), is_currency=True)}
• 52W Low: {format_number(stock_data.get('52_week_low'), is_currency=True)}
• Day Change: {format_number(stock_data.get('day_change_percent'), is_percentage=True)}

💰 FINANCIAL STRENGTH:
• Revenue: {format_number(financial_metrics.get('revenue'), is_currency=True)}
• Market Cap: {format_number(financial_metrics.get('market_cap'), is_currency=True)}
• Profit Margin: {format_number(financial_metrics.get('profit_margin'), is_percentage=True)}
• ROE: {format_number(financial_metrics.get('return_on_equity'), is_percentage=True)}

📊 VALUATION METRICS:
• P/E Ratio: {format_number(financial_metrics.get('pe_ratio'))}
• Price-to-Book: {format_number(financial_metrics.get('price_to_book'))}
• Beta: {format_number(company_data['key_stats'].get('beta'))}

👥 SHAREHOLDERS:
• Shares Outstanding: {format_number(company_data['key_stats'].get('shares_outstanding'))}
• Full-time Employees: {format_number(company_data['key_stats'].get('employees'))}"""

        summary_frame.text = summary_text
        for paragraph in summary_frame.paragraphs:
            paragraph.font.size = Pt(10)
            paragraph.font.color.rgb = RGBColor(52, 73, 94)
        
        # Add major shareholders if available
        if company_data.get('major_shareholders'):
            shareholders_box = slide2.shapes.add_textbox(Inches(1), Inches(6.8), Inches(11), Inches(0.7))
            shareholders_frame = shareholders_box.text_frame
            shareholders_text = "MAJOR SHAREHOLDERS: " + " | ".join([f"{holder.get('Holder', 'N/A')}: {holder.get('Shares', 'N/A')}" for holder in company_data['major_shareholders'][:3]])
            shareholders_frame.text = shareholders_text
            shareholders_frame.paragraphs[0].font.size = Pt(9)
            shareholders_frame.paragraphs[0].font.color.rgb = RGBColor(127, 140, 141)
        
        # Save PowerPoint
        pptx_path = f"/tmp/uploads/{company_name}_analysis.pptx"
        prs.save(pptx_path)
        
        return pptx_path
        
    except Exception as e:
        print(f"Error creating PowerPoint: {e}")
        return None

def extract_ticker_symbol(company_name, exchange=None):
    """
    Try to extract or guess ticker symbol from company name
    """
    try:
        # Simple heuristics for common ticker patterns
        if len(company_name) <= 5 and company_name.isupper():
            return company_name
        
        # Try searching with yfinance
        search_variations = [
            company_name,
            company_name.upper(),
            company_name.replace(' ', ''),
            company_name.replace(' ', '-'),
            company_name.split()[0] if ' ' in company_name else company_name
        ]
        
        for variation in search_variations:
            try:
                # Add exchange suffixes for international markets
                test_symbols = [variation]
                if exchange:
                    if 'TSX' in exchange or 'Toronto' in exchange:
                        test_symbols.append(f"{variation}.TO")
                    elif 'London' in exchange or 'LSE' in exchange:
                        test_symbols.append(f"{variation}.L")
                    elif 'ASX' in exchange or 'Australia' in exchange:
                        test_symbols.append(f"{variation}.AX")
                
                for symbol in test_symbols:
                    try:
                        stock = yf.Ticker(symbol)
                        info = stock.info
                        if info and info.get('longName'):
                            return symbol
                    except:
                        continue
            except:
                continue
        
        return None
    except:
        return None

# For Vercel deployment
if __name__ == "__main__":
    app.run(debug=True) 